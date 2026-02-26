#!/usr/bin/env python3
"""
共立製薬 PPTX合体スクリプト
- スコアシートDX提案 (13枚) → 前半（1.3333倍スケール）
- データ基盤メリット提案 (6枚) → 後半（そのまま）
- 出力スライドサイズ: 12192000 x 6858000 (データ基盤側に統一)
"""

import copy
import os
from lxml import etree
from pptx import Presentation
from pptx.opc.package import Part, _Relationship
from pptx.opc.packuri import PackURI

# --- Paths ---
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
SCORESHEET_PATH = os.path.join(BASE_DIR, "共立製薬-スコアシートDX提案-2026-02.pptx")
DATA_INFRA_PATH = os.path.join(BASE_DIR, "データ基盤メリット提案_共立製薬_NTTDXPN.pptx")
OUTPUT_PATH = os.path.join(BASE_DIR, "共立製薬-統合提案-スコアシートDX+データ基盤_NTTDXPN.pptx")

# --- Constants ---
TARGET_WIDTH = 12192000
TARGET_HEIGHT = 6858000
SOURCE_WIDTH = 9144000
SOURCE_HEIGHT = 5143500
SCALE = TARGET_WIDTH / SOURCE_WIDTH  # 1.3333...

P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


def scale_emu(value):
    return int(round(value * SCALE))


def scale_slide_xml(slide_elem):
    """Scale all shape positions, sizes, font sizes, tables, margins, etc."""
    for xfrm in slide_elem.iter(f'{{{A_NS}}}xfrm'):
        off = xfrm.find(f'{{{A_NS}}}off')
        if off is not None:
            for attr in ('x', 'y'):
                val = off.get(attr)
                if val is not None:
                    off.set(attr, str(scale_emu(int(val))))
        ext = xfrm.find(f'{{{A_NS}}}ext')
        if ext is not None:
            for attr in ('cx', 'cy'):
                val = ext.get(attr)
                if val is not None:
                    ext.set(attr, str(scale_emu(int(val))))

    for tag in ('rPr', 'defRPr', 'endParaRPr'):
        for elem in slide_elem.iter(f'{{{A_NS}}}{tag}'):
            sz = elem.get('sz')
            if sz is not None:
                elem.set('sz', str(int(round(int(sz) * SCALE))))

    for tr in slide_elem.iter(f'{{{A_NS}}}tr'):
        h = tr.get('h')
        if h is not None:
            tr.set('h', str(scale_emu(int(h))))

    for gridCol in slide_elem.iter(f'{{{A_NS}}}gridCol'):
        w = gridCol.get('w')
        if w is not None:
            gridCol.set('w', str(scale_emu(int(w))))

    for pPr in slide_elem.iter(f'{{{A_NS}}}pPr'):
        for attr in ('marL', 'marR', 'indent', 'defTabSz'):
            val = pPr.get(attr)
            if val is not None:
                pPr.set(attr, str(scale_emu(int(val))))

    for bodyPr in slide_elem.iter(f'{{{A_NS}}}bodyPr'):
        for attr in ('lIns', 'tIns', 'rIns', 'bIns'):
            val = bodyPr.get(attr)
            if val is not None:
                bodyPr.set(attr, str(scale_emu(int(val))))

    for ln in slide_elem.iter(f'{{{A_NS}}}ln'):
        w = ln.get('w')
        if w is not None:
            ln.set('w', str(scale_emu(int(w))))

    for tab in slide_elem.iter(f'{{{A_NS}}}tab'):
        pos = tab.get('pos')
        if pos is not None:
            tab.set('pos', str(scale_emu(int(pos))))


def copy_slide_with_images(src_slide, dst_prs, slide_index):
    """
    Copy a slide from source presentation to destination.
    Uses relate_to for images and remaps rIds in XML.
    """
    dst_layout = dst_prs.slide_layouts[0]
    new_slide = dst_prs.slides.add_slide(dst_layout)

    # Build rId mapping: src_rId -> dst_rId
    rid_map = {}

    for rel in src_slide.part.rels.values():
        if 'slideLayout' in rel.reltype or 'notesSlide' in rel.reltype:
            continue

        if rel.is_external:
            new_rid = new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
            rid_map[rel.rId] = new_rid
        else:
            try:
                target_part = rel.target_part
                # Create a new part with unique name
                orig_name = str(target_part.partname)
                parts = orig_name.rsplit('/', 1)
                if len(parts) == 2:
                    new_name = f"{parts[0]}/s{slide_index}_{parts[1]}"
                else:
                    new_name = orig_name

                new_part = Part(
                    PackURI(new_name),
                    target_part.content_type,
                    dst_prs.part.package,
                    target_part.blob,
                )
                new_rid = new_slide.part.relate_to(new_part, rel.reltype)
                rid_map[rel.rId] = new_rid
            except Exception as e:
                print(f"    Warning: Could not copy rel {rel.rId} ({rel.reltype}): {e}")

    # Deep copy source XML
    src_xml = copy.deepcopy(src_slide._element)

    # Remap rIds in the copied XML
    if rid_map:
        xml_str = etree.tostring(src_xml, encoding='unicode')
        for old_rid, new_rid in rid_map.items():
            # Replace r:id="rIdX" and similar references
            xml_str = xml_str.replace(f'"{old_rid}"', f'"{new_rid}"')
        src_xml = etree.fromstring(xml_str)

    # Replace the new slide's content with the source content
    for child in list(new_slide._element):
        new_slide._element.remove(child)
    for attrib_name, attrib_val in src_xml.attrib.items():
        new_slide._element.set(attrib_name, attrib_val)
    for child in src_xml:
        new_slide._element.append(child)

    return new_slide


def main():
    print(f"Loading scoresheet DX: {SCORESHEET_PATH}")
    src_prs = Presentation(SCORESHEET_PATH)
    print(f"  Slides: {len(src_prs.slides)}")
    print(f"  Size: {src_prs.slide_width} x {src_prs.slide_height}")

    print(f"Loading data infra: {DATA_INFRA_PATH}")
    dst_prs = Presentation(DATA_INFRA_PATH)
    print(f"  Slides: {len(dst_prs.slides)}")
    print(f"  Size: {dst_prs.slide_width} x {dst_prs.slide_height}")

    assert src_prs.slide_width == SOURCE_WIDTH, f"Source width mismatch: {src_prs.slide_width}"
    assert dst_prs.slide_width == TARGET_WIDTH, f"Target width mismatch: {dst_prs.slide_width}"

    dst_original_count = len(dst_prs.slides)
    src_slides = list(src_prs.slides)
    print(f"\nCopying {len(src_slides)} slides from scoresheet DX...")

    for i, src_slide in enumerate(src_slides):
        print(f"  Slide {i+1}/{len(src_slides)}...")
        new_slide = copy_slide_with_images(src_slide, dst_prs, i + 1)
        scale_slide_xml(new_slide._element)

    # Reorder: source slides (appended at end) should come first
    print(f"\nReordering slides...")
    src_count = len(src_slides)

    pres_elem = dst_prs._element
    sldIdLst = pres_elem.find(f'{{{P_NS}}}sldIdLst')

    if sldIdLst is not None:
        all_sld_ids = list(sldIdLst)
        src_elems = all_sld_ids[dst_original_count:]
        dst_elems = all_sld_ids[:dst_original_count]

        for elem in list(sldIdLst):
            sldIdLst.remove(elem)
        for elem in src_elems:
            sldIdLst.append(elem)
        for elem in dst_elems:
            sldIdLst.append(elem)

    total = len(dst_prs.slides)
    print(f"\nTotal slides: {total}")
    print(f"  1-{src_count}: ScoreSheet DX (scaled x{SCALE:.4f})")
    print(f"  {src_count+1}-{total}: Data Infra Merit (original)")

    print(f"\nSaving to: {OUTPUT_PATH}")
    dst_prs.save(OUTPUT_PATH)
    print("Done!")


if __name__ == "__main__":
    main()
