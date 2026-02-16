#!/usr/bin/env python3
"""共立製薬 データ基盤構築メリット提案 PPTX生成スクリプト"""

import os
import shutil
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# === 定数 ===
TEMPLATE_PATH = os.path.expanduser("~/.claude/skills/pptx/template.pptx")
LAYOUT_MAIN = 0    # 2_3.メインページ
LAYOUT_TITLE = 1   # 1_タイトルスライド
LAYOUT_BLANK = 2   # 空白
FONT = "Meiryo UI"

# カラーパレット
C_MAIN = RGBColor(0x15, 0x60, 0x82)       # ティールブルー
C_ACCENT = RGBColor(0x6C, 0xE6, 0xE8)     # ライトシアン
C_TEXT = RGBColor(0x33, 0x33, 0x33)        # テキスト濃
C_SUB = RGBColor(0x66, 0x66, 0x66)        # テキスト薄
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)       # 白
C_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)      # 薄グレー
C_BORDER = RGBColor(0xE0, 0xE0, 0xE0)     # ボーダー
C_BLUE_BG = RGBColor(0xED, 0xF7, 0xFA)    # 強調背景
C_GREEN = RGBColor(0x2E, 0x7D, 0x32)      # 成功テキスト
C_GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)   # 成功背景
C_ORANGE = RGBColor(0xE6, 0x7E, 0x22)     # オレンジ
C_GOLD = RGBColor(0xF3, 0xC5, 0x51)       # ゴールド
C_WARN_BG = RGBColor(0xFF, 0xF8, 0xE1)    # 警告背景
C_PILLAR1 = RGBColor(0x15, 0x60, 0x82)    # 柱1: ティールブルー
C_PILLAR2 = RGBColor(0x2E, 0x7D, 0x32)    # 柱2: グリーン


# === ヘルパー関数 ===

def set_font(run, size=None, bold=False, color=None, name=FONT):
    run.font.name = name
    if size:
        run.font.size = Pt(size)
    if bold:
        run.font.bold = True
    if color:
        run.font.color.rgb = color


def set_placeholder_font(placeholder, name=FONT):
    for para in placeholder.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = name


def add_lead_text(slide, text, top=1003303):
    """リード文（★付き）を追加"""
    txBox = slide.shapes.add_textbox(Emu(520700), Emu(top), Emu(11150600), Emu(500000))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "★" + text
    set_font(run, size=16, color=C_SUB)
    return txBox


def add_rounded_box(slide, left, top, width, height, fill_color=None, line_color=None):
    """角丸四角形を追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(left), Emu(top), Emu(width), Emu(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, size=14, bold=False,
                 color=C_TEXT, align=PP_ALIGN.LEFT, line_spacing=None):
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return txBox


def add_number_badge(slide, left, top, text, color=C_MAIN):
    """番号バッジ（円形）を追加"""
    size = 380000
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(left), Emu(top), Emu(size), Emu(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.text = text
    for para in shape.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            set_font(run, size=14, bold=True, color=C_WHITE)
    return shape


# === スライド生成関数 ===

def add_cover_slide(prs):
    """表紙スライド"""
    layout = prs.slide_layouts[LAYOUT_TITLE]
    slide = prs.slides.add_slide(layout)

    slide.placeholders[0].text = "データ基盤構築のメリット"
    set_placeholder_font(slide.placeholders[0])

    # 会社名
    txBox = slide.shapes.add_textbox(Emu(489661), Emu(409009), Emu(3650077), Emu(369332))
    tf = txBox.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "共立製薬株式会社 様"
    set_font(run, size=18)

    # 部署名
    txBox = slide.shapes.add_textbox(Emu(788919), Emu(4725384), Emu(7249488), Emu(369332))
    tf = txBox.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "株式会社NTT DXパートナー"
    set_font(run, size=18)

    # 日付
    txBox = slide.shapes.add_textbox(Emu(788919), Emu(4356052), Emu(2205681), Emu(369332))
    tf = txBox.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "2026年2月"
    set_font(run, size=18)

    return slide


def add_overview_slide(prs):
    """全体像スライド - 2つの柱を視覚的に表現"""
    layout = prs.slide_layouts[LAYOUT_MAIN]
    slide = prs.slides.add_slide(layout)

    slide.placeholders[0].text = "データ基盤構築のメリット 全体像"
    set_placeholder_font(slide.placeholders[0])

    add_lead_text(slide, "データ基盤を構築することで、研究の質向上と業務効率化の両面でメリットが得られます。")

    # 中央メッセージ
    msg_top = 1550000
    add_text_box(slide, 520700, msg_top, 11150600, 400000,
                 "経験や勘ではなく、データに基づいた意思決定が可能になる",
                 size=16, bold=True, color=C_MAIN, align=PP_ALIGN.CENTER)

    # === 柱1: データを研究に活かせる ===
    p1_left = 520700
    p1_top = 2200000
    p1_w = 5200000
    p1_h = 3800000

    # 柱1 背景
    add_rounded_box(slide, p1_left, p1_top, p1_w, p1_h, fill_color=C_WHITE, line_color=C_PILLAR1)
    # 柱1 ヘッダーバー
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(p1_left), Emu(p1_top), Emu(p1_w), Emu(500000)
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = C_PILLAR1
    hdr.line.fill.background()
    hdr.text = "1. データを研究に活かせる"
    for para in hdr.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            set_font(run, size=18, bold=True, color=C_WHITE)

    # 1-1
    sub_top = p1_top + 650000
    add_number_badge(slide, p1_left + 200000, sub_top, "1-1", C_PILLAR1)
    add_text_box(slide, p1_left + 680000, sub_top + 30000, 4200000, 350000,
                 "傾向の可視化", size=16, bold=True, color=C_TEXT)
    add_text_box(slide, p1_left + 680000, sub_top + 380000, 4200000, 600000,
                 "ロット間比較、観察者間のばらつき分析など、\nデータの傾向を定量的に把握できる",
                 size=12, color=C_SUB, line_spacing=20)

    # 1-2
    sub_top2 = sub_top + 1200000
    add_number_badge(slide, p1_left + 200000, sub_top2, "1-2", C_PILLAR1)
    add_text_box(slide, p1_left + 680000, sub_top2 + 30000, 4200000, 350000,
                 "異常検出・示唆出し", size=16, bold=True, color=C_TEXT)
    add_text_box(slide, p1_left + 680000, sub_top2 + 380000, 4200000, 600000,
                 "早期予測、適正頭数算出、\nHumane Endpoint基準の客観化",
                 size=12, color=C_SUB, line_spacing=20)

    # === 柱2: 業務を効率化できる ===
    p2_left = 6000000
    p2_top = 2200000
    p2_w = 5700000
    p2_h = 3800000

    # 柱2 背景
    add_rounded_box(slide, p2_left, p2_top, p2_w, p2_h, fill_color=C_WHITE, line_color=C_PILLAR2)
    # 柱2 ヘッダーバー
    hdr2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(p2_left), Emu(p2_top), Emu(p2_w), Emu(500000)
    )
    hdr2.fill.solid()
    hdr2.fill.fore_color.rgb = C_PILLAR2
    hdr2.line.fill.background()
    hdr2.text = "2. 業務を効率化できる"
    for para in hdr2.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            set_font(run, size=18, bold=True, color=C_WHITE)

    # 2-1 ~ 2-4
    items = [
        ("2-1", "過去データの検索"),
        ("2-2", "転記・格納作業の削減"),
        ("2-3", "規制対応（GCP査察）"),
        ("2-4", "申請資料の作成"),
    ]
    for i, (num, label) in enumerate(items):
        y = p2_top + 650000 + i * 700000
        add_number_badge(slide, p2_left + 200000, y, num, C_PILLAR2)
        add_text_box(slide, p2_left + 680000, y + 50000, 4600000, 350000,
                     label, size=14, bold=True, color=C_TEXT)

    return slide


def add_pillar1_detail_slide(prs):
    """柱1: データを研究に活かせる 詳細"""
    layout = prs.slide_layouts[LAYOUT_MAIN]
    slide = prs.slides.add_slide(layout)

    slide.placeholders[0].text = "1. データを研究に活かせる"
    set_placeholder_font(slide.placeholders[0])

    add_lead_text(slide, "蓄積されたデータを分析することで、傾向の把握や異常の早期検出が可能になります。")

    # === 1-1 傾向の可視化 カード ===
    card1_left = 520700
    card1_top = 1700000
    card1_w = 5400000
    card1_h = 4200000

    add_rounded_box(slide, card1_left, card1_top, card1_w, card1_h,
                    fill_color=C_WHITE, line_color=C_PILLAR1)
    # ヘッダー
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(card1_left), Emu(card1_top), Emu(card1_w), Emu(450000)
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = C_PILLAR1
    hdr.line.fill.background()
    hdr.text = "1-1. 傾向の可視化"
    for para in hdr.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            set_font(run, size=16, bold=True, color=C_WHITE)

    # シーン1
    s1_top = card1_top + 600000
    add_text_box(slide, card1_left + 200000, s1_top, 5000000, 300000,
                 "ロット間で品質に差がないか確認したい", size=13, bold=True, color=C_TEXT)
    # Before/After
    add_rounded_box(slide, card1_left + 200000, s1_top + 350000, 2400000, 700000,
                    fill_color=C_WARN_BG, line_color=C_GOLD)
    add_text_box(slide, card1_left + 300000, s1_top + 380000, 2200000, 250000,
                 "現在（想定）", size=10, bold=True, color=C_ORANGE)
    add_text_box(slide, card1_left + 300000, s1_top + 600000, 2200000, 400000,
                 "手作業で突き合わせ比較\n担当者の判断に依存", size=10, color=C_SUB, line_spacing=16)

    # 矢印
    add_text_box(slide, card1_left + 2650000, s1_top + 550000, 350000, 300000,
                 "→", size=20, bold=True, color=C_MAIN, align=PP_ALIGN.CENTER)

    add_rounded_box(slide, card1_left + 3000000, s1_top + 350000, 2200000, 700000,
                    fill_color=C_GREEN_BG, line_color=C_PILLAR2)
    add_text_box(slide, card1_left + 3100000, s1_top + 380000, 2000000, 250000,
                 "データ基盤構築後", size=10, bold=True, color=C_GREEN)
    add_text_box(slide, card1_left + 3100000, s1_top + 600000, 2000000, 400000,
                 "Golden Batchと自動比較\n逸脱を即座に検出", size=10, color=C_SUB, line_spacing=16)

    # シーン2
    s2_top = s1_top + 1250000
    add_text_box(slide, card1_left + 200000, s2_top, 5000000, 300000,
                 "観察者間のスコアばらつきを確認したい", size=13, bold=True, color=C_TEXT)
    add_rounded_box(slide, card1_left + 200000, s2_top + 350000, 2400000, 700000,
                    fill_color=C_WARN_BG, line_color=C_GOLD)
    add_text_box(slide, card1_left + 300000, s2_top + 380000, 2200000, 250000,
                 "現在（想定）", size=10, bold=True, color=C_ORANGE)
    add_text_box(slide, card1_left + 300000, s2_top + 600000, 2200000, 400000,
                 "研修や口頭指導で統一", size=10, color=C_SUB, line_spacing=16)

    add_text_box(slide, card1_left + 2650000, s2_top + 550000, 350000, 300000,
                 "→", size=20, bold=True, color=C_MAIN, align=PP_ALIGN.CENTER)

    add_rounded_box(slide, card1_left + 3000000, s2_top + 350000, 2200000, 700000,
                    fill_color=C_GREEN_BG, line_color=C_PILLAR2)
    add_text_box(slide, card1_left + 3100000, s2_top + 380000, 2000000, 250000,
                 "データ基盤構築後", size=10, bold=True, color=C_GREEN)
    add_text_box(slide, card1_left + 3100000, s2_top + 600000, 2000000, 400000,
                 "統計的に比較\n評価傾向の差を定量化", size=10, color=C_SUB, line_spacing=16)

    # === 1-2 異常検出・示唆出し カード ===
    card2_left = 6200000
    card2_top = 1700000
    card2_w = 5500000
    card2_h = 4200000

    add_rounded_box(slide, card2_left, card2_top, card2_w, card2_h,
                    fill_color=C_WHITE, line_color=C_PILLAR1)
    hdr2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(card2_left), Emu(card2_top), Emu(card2_w), Emu(450000)
    )
    hdr2.fill.solid()
    hdr2.fill.fore_color.rgb = C_PILLAR1
    hdr2.line.fill.background()
    hdr2.text = "1-2. 異常検出・示唆出し"
    for para in hdr2.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            set_font(run, size=16, bold=True, color=C_WHITE)

    # シーン1: 早期予測
    s1_top = card2_top + 600000
    add_text_box(slide, card2_left + 200000, s1_top, 5100000, 300000,
                 "試験途中から結果を早期に予測したい", size=13, bold=True, color=C_TEXT)
    add_rounded_box(slide, card2_left + 200000, s1_top + 350000, 2400000, 550000,
                    fill_color=C_WARN_BG, line_color=C_GOLD)
    add_text_box(slide, card2_left + 300000, s1_top + 380000, 2200000, 200000,
                 "現在（想定）", size=10, bold=True, color=C_ORANGE)
    add_text_box(slide, card2_left + 300000, s1_top + 560000, 2200000, 300000,
                 "14日間完了後に評価", size=10, color=C_SUB)

    add_text_box(slide, card2_left + 2650000, s1_top + 450000, 350000, 300000,
                 "→", size=20, bold=True, color=C_MAIN, align=PP_ALIGN.CENTER)

    add_rounded_box(slide, card2_left + 3000000, s1_top + 350000, 2300000, 550000,
                    fill_color=C_GREEN_BG, line_color=C_PILLAR2)
    add_text_box(slide, card2_left + 3100000, s1_top + 380000, 2100000, 200000,
                 "データ基盤構築後", size=10, bold=True, color=C_GREEN)
    add_text_box(slide, card2_left + 3100000, s1_top + 560000, 2100000, 300000,
                 "蓄積データから早期判断", size=10, color=C_SUB)

    # シーン2: 適正頭数
    s2_top = s1_top + 1050000
    add_text_box(slide, card2_left + 200000, s2_top, 5100000, 300000,
                 "試験に使う動物の頭数を適切に決めたい", size=13, bold=True, color=C_TEXT)
    add_rounded_box(slide, card2_left + 200000, s2_top + 350000, 2400000, 550000,
                    fill_color=C_WARN_BG, line_color=C_GOLD)
    add_text_box(slide, card2_left + 300000, s2_top + 380000, 2200000, 200000,
                 "現在（想定）", size=10, bold=True, color=C_ORANGE)
    add_text_box(slide, card2_left + 300000, s2_top + 560000, 2200000, 300000,
                 "慣例・経験値で決定", size=10, color=C_SUB)

    add_text_box(slide, card2_left + 2650000, s2_top + 450000, 350000, 300000,
                 "→", size=20, bold=True, color=C_MAIN, align=PP_ALIGN.CENTER)

    add_rounded_box(slide, card2_left + 3000000, s2_top + 350000, 2300000, 550000,
                    fill_color=C_GREEN_BG, line_color=C_PILLAR2)
    add_text_box(slide, card2_left + 3100000, s2_top + 380000, 2100000, 200000,
                 "データ基盤構築後", size=10, bold=True, color=C_GREEN)
    add_text_box(slide, card2_left + 3100000, s2_top + 560000, 2100000, 300000,
                 "統計的に最小頭数を算出\n3Rs(Reduction)に貢献", size=10, color=C_SUB, line_spacing=16)

    # シーン3: Humane Endpoint
    s3_top = s2_top + 1050000
    add_text_box(slide, card2_left + 200000, s3_top, 5100000, 300000,
                 "Humane Endpoint基準を客観化したい", size=13, bold=True, color=C_TEXT)
    add_rounded_box(slide, card2_left + 200000, s3_top + 350000, 2400000, 550000,
                    fill_color=C_WARN_BG, line_color=C_GOLD)
    add_text_box(slide, card2_left + 300000, s3_top + 380000, 2200000, 200000,
                 "現在（想定）", size=10, bold=True, color=C_ORANGE)
    add_text_box(slide, card2_left + 300000, s3_top + 560000, 2200000, 300000,
                 "経験則に基づく判断", size=10, color=C_SUB)

    add_text_box(slide, card2_left + 2650000, s3_top + 450000, 350000, 300000,
                 "→", size=20, bold=True, color=C_MAIN, align=PP_ALIGN.CENTER)

    add_rounded_box(slide, card2_left + 3000000, s3_top + 350000, 2300000, 550000,
                    fill_color=C_GREEN_BG, line_color=C_PILLAR2)
    add_text_box(slide, card2_left + 3100000, s3_top + 380000, 2100000, 200000,
                 "データ基盤構築後", size=10, bold=True, color=C_GREEN)
    add_text_box(slide, card2_left + 3100000, s3_top + 560000, 2100000, 300000,
                 "エビデンスで倫理委員会\nへの説明根拠を提示", size=10, color=C_SUB, line_spacing=16)

    return slide


def add_pillar2_detail_slide(prs):
    """柱2: 業務を効率化できる 詳細"""
    layout = prs.slide_layouts[LAYOUT_MAIN]
    slide = prs.slides.add_slide(layout)

    slide.placeholders[0].text = "2. 業務を効率化できる"
    set_placeholder_font(slide.placeholders[0])

    add_lead_text(slide, "データの検索・転記・保管にかかる手作業を削減し、規制対応や申請業務を効率化します。")

    # 4つのカードを2x2グリッドで配置
    items = [
        {
            "num": "2-1", "title": "過去データの検索",
            "scene": "過去の試験結果を調べたい",
            "before": "書庫から紙を探索\n担当者の異動で所在不明リスク",
            "after": "試験名・個体番号等で\n即座に横断検索",
        },
        {
            "num": "2-2", "title": "転記・格納作業の削減",
            "scene": "スコアシートのデータを記録・保管したい",
            "before": "手書き → ファイリング\n他帳票への転記作業が発生",
            "after": "タブレットからデジタル入力\nDBに自動格納・転記ミス解消",
        },
        {
            "num": "2-3", "title": "規制対応（GCP査察）",
            "scene": "GCP査察・監査に対応したい",
            "before": "紙の記録を手作業で整理\n修正履歴の追跡が困難",
            "after": "デジタル監査証跡が自動\nraw data即時提出可能",
        },
        {
            "num": "2-4", "title": "申請資料の作成",
            "scene": "承認申請用の資料を作成したい",
            "before": "集計・整形を手作業で実施",
            "after": "統計解析・個体別データを\n申請様式で自動出力",
        },
    ]

    card_w = 5400000
    card_h = 2500000
    gap_x = 500000
    gap_y = 300000
    start_left = 520700
    start_top = 1700000

    for i, item in enumerate(items):
        col = i % 2
        row = i // 2
        left = start_left + col * (card_w + gap_x)
        top = start_top + row * (card_h + gap_y)

        # カード背景
        add_rounded_box(slide, left, top, card_w, card_h, fill_color=C_WHITE, line_color=C_PILLAR2)

        # 番号バッジ + タイトル
        add_number_badge(slide, left + 150000, top + 100000, item["num"], C_PILLAR2)
        add_text_box(slide, left + 600000, top + 120000, 4500000, 350000,
                     item["title"], size=15, bold=True, color=C_TEXT)

        # シーン
        add_text_box(slide, left + 200000, top + 500000, 5000000, 250000,
                     item["scene"], size=11, bold=True, color=C_SUB)

        # Before
        bf_top = top + 800000
        add_rounded_box(slide, left + 200000, bf_top, 2300000, 700000,
                        fill_color=C_WARN_BG, line_color=C_GOLD)
        add_text_box(slide, left + 300000, bf_top + 30000, 2100000, 200000,
                     "現在（想定）", size=9, bold=True, color=C_ORANGE)
        add_text_box(slide, left + 300000, bf_top + 230000, 2100000, 430000,
                     item["before"], size=9, color=C_SUB, line_spacing=15)

        # 矢印
        add_text_box(slide, left + 2550000, bf_top + 200000, 300000, 300000,
                     "→", size=18, bold=True, color=C_PILLAR2, align=PP_ALIGN.CENTER)

        # After
        add_rounded_box(slide, left + 2850000, bf_top, 2350000, 700000,
                        fill_color=C_GREEN_BG, line_color=C_PILLAR2)
        add_text_box(slide, left + 2950000, bf_top + 30000, 2150000, 200000,
                     "データ基盤構築後", size=9, bold=True, color=C_GREEN)
        add_text_box(slide, left + 2950000, bf_top + 230000, 2150000, 430000,
                     item["after"], size=9, color=C_SUB, line_spacing=15)

    return slide


def add_detail_table_slide(prs):
    """活用シーン一覧テーブル"""
    layout = prs.slide_layouts[LAYOUT_MAIN]
    slide = prs.slides.add_slide(layout)

    slide.placeholders[0].text = "具体的な活用シーン一覧"
    set_placeholder_font(slide.placeholders[0])

    add_lead_text(slide, "各メリットに紐づく具体的な活用シーンと、データ基盤構築前後の変化を整理しています。")

    # テーブル
    headers = ["分類", "活用シーン", "現在の運用（想定）", "データ基盤構築後"]
    rows_data = [
        ["1-1", "ロット間品質比較", "手作業で突き合わせ比較", "Golden Batchと自動比較"],
        ["1-1", "観察者間ばらつき確認", "研修・口頭指導で統一", "統計的に評価傾向を定量化"],
        ["1-2", "試験結果の早期予測", "14日間完了後に評価", "蓄積データから早期判断"],
        ["1-2", "適正頭数の決定", "慣例・経験値で決定", "統計的に最小頭数を算出"],
        ["1-2", "Humane Endpoint基準", "経験則に基づく判断", "エビデンスで説明根拠を提示"],
        ["2-1", "過去試験結果の検索", "書庫から紙を探索", "キーワードで即座に横断検索"],
        ["2-2", "スコアシート記録・保管", "手書き→ファイリング", "デジタル入力→DB自動格納"],
        ["2-3", "GCP査察・監査対応", "紙記録を手作業で整理", "デジタル監査証跡が自動記録"],
        ["2-4", "承認申請資料作成", "集計・整形を手作業", "申請様式で自動出力"],
    ]

    num_rows = len(rows_data) + 1
    num_cols = len(headers)
    left = Emu(300000)
    top = Emu(1700000)
    width = Emu(11600000)
    height = Emu(380000 * num_rows)

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # 列幅
    table.columns[0].width = Emu(800000)
    table.columns[1].width = Emu(2600000)
    table.columns[2].width = Emu(4100000)
    table.columns[3].width = Emu(4100000)

    # ヘッダー
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                set_font(run, size=11, bold=True)

    # データ
    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                if ci == 0:
                    para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    set_font(run, size=10)
                    if ci == 0:
                        run.font.bold = True
                        run.font.color.rgb = C_MAIN

    # 注釈
    add_text_box(slide, 300000, top.emu + height.emu + 100000, 11600000, 300000,
                 "※「現在の運用（想定）」列は推測であり、実際の業務フローはヒアリングによる確認が必要です。",
                 size=10, color=C_SUB)

    return slide


def add_summary_slide(prs):
    """まとめスライド"""
    layout = prs.slide_layouts[LAYOUT_MAIN]
    slide = prs.slides.add_slide(layout)

    slide.placeholders[0].text = "まとめ"
    set_placeholder_font(slide.placeholders[0])

    add_lead_text(slide, "データ基盤構築により、研究品質の向上と業務効率化を同時に実現できます。")

    # 中央のまとめボックス
    box_left = 1500000
    box_top = 1800000
    box_w = 9200000
    box_h = 3600000

    add_rounded_box(slide, box_left, box_top, box_w, box_h, fill_color=C_BLUE_BG, line_color=C_MAIN)

    # 柱1
    add_rounded_box(slide, box_left + 300000, box_top + 300000, 4000000, 450000,
                    fill_color=C_PILLAR1)
    add_text_box(slide, box_left + 300000, box_top + 380000, 4000000, 300000,
                 "1. データを研究に活かせる", size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, box_left + 500000, box_top + 850000, 3600000, 300000,
                 "✓ 傾向の可視化（ロット比較・ばらつき分析）", size=12, color=C_TEXT)
    add_text_box(slide, box_left + 500000, box_top + 1150000, 3600000, 300000,
                 "✓ 異常検出・示唆出し（早期予測・適正頭数）", size=12, color=C_TEXT)

    # 柱2
    add_rounded_box(slide, box_left + 4700000, box_top + 300000, 4200000, 450000,
                    fill_color=C_PILLAR2)
    add_text_box(slide, box_left + 4700000, box_top + 380000, 4200000, 300000,
                 "2. 業務を効率化できる", size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, box_left + 4900000, box_top + 850000, 3800000, 300000,
                 "✓ 過去データの即時検索", size=12, color=C_TEXT)
    add_text_box(slide, box_left + 4900000, box_top + 1100000, 3800000, 300000,
                 "✓ 転記・格納作業の削減", size=12, color=C_TEXT)
    add_text_box(slide, box_left + 4900000, box_top + 1350000, 3800000, 300000,
                 "✓ GCP査察・監査証跡の自動化", size=12, color=C_TEXT)
    add_text_box(slide, box_left + 4900000, box_top + 1600000, 3800000, 300000,
                 "✓ 申請資料の自動出力", size=12, color=C_TEXT)

    # 底部メッセージ
    add_rounded_box(slide, box_left + 1500000, box_top + 2200000, 6200000, 600000,
                    fill_color=C_MAIN)
    add_text_box(slide, box_left + 1500000, box_top + 2300000, 6200000, 400000,
                 "経験や勘ではなく、データに基づいた意思決定へ",
                 size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 注釈
    add_text_box(slide, 520700, 5800000, 11150600, 300000,
                 "※「現在の運用（想定）」部分は推測です。実際の業務フローはヒアリングにて確認させていただきます。",
                 size=10, color=C_SUB)

    return slide


def delete_template_slides(prs, count=5):
    """テンプレートの最初のcount枚のスライドを削除"""
    for _ in range(count):
        rId = None
        rId_attr = prs.slides._sldIdLst[0].attrib
        for key in rId_attr:
            if key.endswith('}id'):
                rId = rId_attr[key]
                break
        if rId:
            prs.part.drop_rel(rId)
        sldId = prs.slides._sldIdLst[0]
        prs.slides._sldIdLst.remove(sldId)


# === メイン ===

def main():
    # テンプレートをコピーして開く
    working_copy = tempfile.mktemp(suffix=".pptx")
    shutil.copy2(TEMPLATE_PATH, working_copy)
    prs = Presentation(working_copy)

    # スライド生成
    add_cover_slide(prs)           # 1. 表紙
    add_overview_slide(prs)        # 2. メリット全体像
    add_pillar1_detail_slide(prs)  # 3. 柱1 詳細
    add_pillar2_detail_slide(prs)  # 4. 柱2 詳細
    add_detail_table_slide(prs)    # 5. 活用シーン一覧
    add_summary_slide(prs)         # 6. まとめ

    # テンプレートスライドを削除
    delete_template_slides(prs, 5)

    # 保存
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, "データ基盤メリット提案_共立製薬_NTTDXPN.pptx")
    prs.save(output_path)

    # 作業用コピーを削除
    os.remove(working_copy)

    print(f"Saved: {output_path}")


if __name__ == "__main__":
    main()
